# -*- coding: utf-8 -*-
"""
Medici Ailesi — 8 Dakikalık Sunum (Arial, açık tema, görselli)
Gereksinimler:
  - Python 3.8+
  - pip install python-pptx requests

Yerelde çalıştırmak için:
  python build_medici_presentation.py
Çıktı:
  Medici_Ailesi_Sunum.pptx (16:9)

Not: Betik, Wikimedia Commons'tan bazı kamu malı / serbest lisanslı görselleri indirir.
URL'ler çalışmazsa CREDITS.md'deki bağlantıları güncelleyebilirsiniz.
"""

import os
import io
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

ASSETS_DIR = "assets"

def ensure_dir(path):
    if not os.path.exists(path):
        os.makedirs(path, exist_ok=True)

def download_image(url, dest_path, timeout=30):
    try:
        r = requests.get(url, timeout=timeout)
        r.raise_for_status()
        with open(dest_path, "wb") as f:
            f.write(r.content)
        return True
    except Exception as e:
        print(f"[WARN] Görsel indirilemedi: {url} -> {e}")
        return False

def set_textframe_arial(text_frame, font_size_pt=None, align_left=True):
    for p in text_frame.paragraphs:
        if align_left:
            p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.name = "Arial"
            if font_size_pt:
                r.font.size = Pt(font_size_pt)

def add_bullet_slide(prs, title, bullets, notes, images=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    # Başlık
    slide.shapes.title.text = title
    set_textframe_arial(slide.shapes.title.text_frame, font_size_pt=36, align_left=False)

    # İçerik maddeleri
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            body.text = b
        else:
            p = body.add_paragraph()
            p.text = b
            p.level = 0
    set_textframe_arial(body, font_size_pt=22)

    # Konuşmacı notları
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = notes.strip()
    set_textframe_arial(notes_tf, font_size_pt=14)

    # Görseller
    if images:
        for img in images:
            url = img.get("url")
            filename = img.get("filename")
            left = Inches(img.get("left_in", 7.0))
            top = Inches(img.get("top_in", 1.2))
            width_in = img.get("width_in")
            height_in = img.get("height_in")

            local_path = os.path.join(ASSETS_DIR, filename)
            if url and not os.path.exists(local_path):
                ok = download_image(url, local_path)
                if not ok:
                    continue
            if os.path.exists(local_path):
                if width_in and height_in:
                    slide.shapes.add_picture(local_path, left, top, width=Inches(width_in), height=Inches(height_in))
                elif width_in:
                    slide.shapes.add_picture(local_path, left, top, width=Inches(width_in))
                elif height_in:
                    slide.shapes.add_picture(local_path, left, top, height=Inches(height_in))
                else:
                    slide.shapes.add_picture(local_path, left, top, width=Inches(5.8))
    return slide

def main():
    ensure_dir(ASSETS_DIR)
    prs = Presentation()
    # 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # SLIDE 1 — Başlık
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Medici Ailesi — Bankadan Şehre"
    set_textframe_arial(s1.shapes.title.text_frame, font_size_pt=40, align_left=False)
    s1.placeholders[1].text = "Finans, Sanat ve Kubbe — 8 Dakikalık Sunum\nFinans → Sanat → Şehir"
    set_textframe_arial(s1.placeholders[1].text_frame, font_size_pt=22, align_left=False)
    n1 = ("Bugün, bankacılığın imgeye; imgenin de kalıcı bir şehir kimliğine nasıl dönüştüğünü konuşacağız.\n"
          "Kahramanımız Medici Ailesi; sahnemiz 15. yy Floransa’sı. Ana fikir: Sermaye görünür yatırıma dönüşürse sanat olur; "
          "sanat şehrin hafızasına kazınırsa kimlik olur.")
    s1.notes_slide.notes_text_frame.text = n1
    set_textframe_arial(s1.notes_slide.notes_text_frame, font_size_pt=14)

    # Görsel (harita) — isteğe bağlı
    img1_url = "https://upload.wikimedia.org/wikipedia/commons/0/0c/Stefano_Buonsignori_-_Novae_urbs_Florentiae_descriptio_-_1584.jpg"
    img1_path = os.path.join(ASSETS_DIR, "florence_map_1584.jpg")
    if download_image(img1_url, img1_path):
        s1.shapes.add_picture(img1_path, Inches(7.2), Inches(1.5), width=Inches(5.6))

    # SLIDE 2 — Floransa’da Zemin
    add_bullet_slide(
        prs,
        "Floransa’da Zemin: İklim ve Kurumlar",
        [
            "Cumhuriyet Floransa’sı: rekabet ve itibar kültürü",
            "Loncalar (Yün–İpek): standart, yarışma, bütçe",
            "Ticaret ağları: Akdeniz ⟷ Kuzey Avrupa"
        ],
        notes=(
            "Floransa, cumhuriyet kurumları ve itibar rekabetiyle dinamikti. Loncalar yalnız üretimi değil, kamusal projeleri ve yarışmaları da yönetiyordu. "
            "Akdeniz’den Kuzey Avrupa’ya sermaye ve bilgi akışı kente yöneldi; finansal inovasyon kök saldı."
        ),
        images=[
            {
                "url": "https://upload.wikimedia.org/wikipedia/commons/0/0c/Stefano_Buonsignori_-_Novae_urbs_Florentiae_descriptio_-_1584.jpg",
                "filename": "florence_map_1584.jpg",
                "left_in": 7.1, "top_in": 1.2, "width_in": 5.7
            }
        ]
    )

    # SLIDE 3 — Medici Bankası
    add_bullet_slide(
        prs,
        "Medici Bankası: Güvenin Mimarisi",
        [
            "1397: Giovanni di Bicci — ölçeklenebilir model",
            "Papalık bankacılığı: yüksek güven, görünürlük",
            "Şube ağı: Roma–Venedik–Londra–Brugge",
            "Formül: Güven + Ağ = Ölçek"
        ],
        notes=(
            "Disiplinli defter tutma, temkinli risk ve temsilcilikler Medici’yi Avrupa finansının omurgasına yerleştirdi. "
            "Banka, yalnız finansman değil, şehir için stratejik bir araçtı."
        ),
        images=[
            {
                "url": "https://upload.wikimedia.org/wikipedia/commons/9/9e/Pacioli_Summa_1494.jpg",
                "filename": "pacioli_summa_1494.jpg",
                "left_in": 7.1, "top_in": 1.2, "width_in": 5.7
            }
        ]
    )

    # SLIDE 4 — Neden Sanat?
    add_bullet_slide(
        prs,
        "Neden Sanat? Meşruiyet, Propaganda, Hafıza",
        [
            "Görünür iyilik: özel servetten kamusal fayda",
            "Kamusal–özel sınırını bulanıklaştırma",
            "Kent markası: “Medici” = Floransa imgesi"
        ],
        notes=(
            "Sanat, gücün en kalıcı dili. Medici, saray–şapel–şenlik gibi araçlarla adı kentsel hafızaya kazıdı; "
            "hami yalnız para veren değil, zevki ve ideali şekillendiren aktör oldu."
        ),
        images=[
            {
                "url": "https://upload.wikimedia.org/wikipedia/commons/2/2e/Palazzo_Medici_Riccardi_-_Firenze.JPG",
                "filename": "palazzo_medici_riccardi.jpg",
                "left_in": 7.1, "top_in": 1.2, "width_in": 5.7
            }
        ]
    )

    # SLIDE 5 — Yaratıcı Ekosistem
    add_bullet_slide(
        prs,
        "Yaratıcı Ekosistem: Atölyeler ve Ağlar",
        [
            "Botticelli: Primavera, Venüs’ün Doğuşu",
            "Verrocchio/Ghirlandaio: usta-çırak-sipariş ağı",
            "Leonardo: Floransa yılları ve etkileşim"
        ],
        notes=(
            "Medici çevresi, yeteneği erken yakalayıp doğru bağlamla buluşturdu. "
            "Sipariş ve fikir akışı hızlandı; mitoloji ve neoplatonizmle bezenmiş bir dil doğdu."
        ),
        images=[
            {
                "url": "https://upload.wikimedia.org/wikipedia/commons/3/3c/Sandro_Botticelli_-_Primavera_-_Google_Art_Project.jpg",
                "filename": "botticelli_primavera.jpg",
                "left_in": 7.1, "top_in": 0.9, "width_in": 5.7
            }
        ]
    )

    # SLIDE 6 — Michelangelo ve Medici
    add_bullet_slide(
        prs,
        "Michelangelo: Lorenzo’nun Bahçesinden Şapele",
        [
            "Lorenzo il Magnifico: heykel bahçesi, erken himaye",
            "Medici Şapeli (Yeni Sakristi): Gece–Gündüz",
            "Mesaj: Zaman, güç, fanilik — mermerde felsefe"
        ],
        notes=(
            "Lorenzo’nun himayesi, Michelangelo’nun formunu erken dönemde şekillendirdi. "
            "Yeni Sakristi’deki alegoriler aile anıtının ötesinde bir düşünce metnidir."
        ),
        images=[
            {
                "url": "https://upload.wikimedia.org/wikipedia/commons/4/42/Sagrestia_Nuova_de_Michelangelo.jpg",
                "filename": "medici_chapel_interior.jpg",
                "left_in": 7.1, "top_in": 0.9, "width_in": 5.7
            },
            {
                "url": "https://upload.wikimedia.org/wikipedia/commons/d/de/Night_by_Michelangelo_JBU001.jpg",
                "filename": "michelangelo_night.jpg",
                "left_in": 7.1, "top_in": 4.2, "width_in": 2.8
            }
        ]
    )

    # SLIDE 7 — Duomo Kubbesi
    add_bullet_slide(
        prs,
        "Duomo Kubbesi: “İmkânsız”ın Mühendisliği",
        [
            "Sorun: Dev açıklık, merkezleme iskelesi yok",
            "Brunelleschi: çift kabuk, balıkkılçığı tuğla örgüsü",
            "Özgün makineler: vinç/dişli/halat sistemleri",
            "Cesaretin ödüllendirilmesi: lonca + şehir iklimi"
        ],
        notes=(
            "Brunelleschi, çift kabuk ve balıkkılçığı örgü ile kuvvetleri yönetti; özgün makinelerle lojistiği çözdü. "
            "Teknik yenilik, kültürel güven olmadan yeşermez; bu iklimi Medici etkisi kurdu."
        ),
        images=[
            {
                # 1911 Encyclopaedia Britannica - Section diagram (PD)
                "url": "https://upload.wikimedia.org/wikipedia/commons/3/35/EB1911_-_Dome_-_Fig._8.—Section_through_the_Dome_of_Florence._%28From_Chapuy%27s_Monde_Classique%29.png",
                "filename": "duomo_section_eb1911.png",
                "left_in": 7.1, "top_in": 0.9, "width_in": 5.7
            }
        ]
    )

    # SLIDE 8 — Modelin Formülü
    add_bullet_slide(
        prs,
        "Modelin Formülü: Bankadan Şehre",
        [
            "Sermaye → Patronaj → Kentsel miras",
            "“Görünür yatırım” kalıcılık üretir",
            "Ders: güven, ağ, vizyon, imge birlikteliği"
        ],
        notes=(
            "Sermaye güven ve ağlarla büyür; görünür yatırıma dönünce sanat olur; sanat hafızaya kazınınca kimlik üretir. "
            "Bugünün kurumları için soru: Hangi yatırımınız kamusal hafızada kalıcı?"
        ),
        images=[]
    )

    # SLIDE 9 — Kapanış
    add_bullet_slide(
        prs,
        "Kapanış ve Soru",
        [
            "“Kaynağı büyüt, imgeye çevir, şehre kazı.”",
            "Teşekkürler — Sorular?"
        ],
        notes=(
            "Özet: Medici modeli üç adımda. İsterseniz kubbe tekniğinin kısa bir şemasıyla başlayabiliriz."
        ),
        images=[]
    )

    prs.save("Medici_Ailesi_Sunum.pptx")
    print("Medici_Ailesi_Sunum.pptx oluşturuldu")

if __name__ == "__main__":
    main()
